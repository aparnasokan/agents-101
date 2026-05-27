"""
rag.py
Local multi-source RAG store for Agent 101.

Supports:
- Multiple source folders/files via RAG_SOURCE_DIRS separated by semicolons
- Single source fallback via RAG_SOURCE_DIR
- HTML/HTM, PPTX, JSON, TXT, MD ingestion
- Azure OpenAI embeddings
- FAISS local vector index
- Debug logs for source discovery, parsing, indexing, and retrieval

Required env:
  AZURE_OPENAI_API_KEY=...
  AZURE_OPENAI_ENDPOINT=...
  AZURE_OPENAI_EMBEDDING_DEPLOYMENT=your-embedding-deployment

Optional env:
  RAG_SOURCE_DIRS=/mnt/c/path/to/use-cases;/home/user/app/knowledge
  RAG_SOURCE_DIR=knowledge
  RAG_INDEX_DIR=.rag_index
  RAG_FORCE_REBUILD=true
  RAG_DEBUG=true
"""

from __future__ import annotations

import hashlib
import json
import os
import pickle
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any

from dotenv import load_dotenv

load_dotenv()

print("[RAG] Loaded complete multi-source rag.py")

try:
    import faiss  # type: ignore
except Exception:  # pragma: no cover
    faiss = None

try:
    from bs4 import BeautifulSoup  # type: ignore
except Exception:  # pragma: no cover
    BeautifulSoup = None

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None

from openai import AzureOpenAI


ALLOWED_EXTENSIONS = {".html", ".htm", ".pptx", ".json", ".txt", ".md"}
DEFAULT_INDEX_DIR = ".rag_index"


def _first_env(*names: str) -> str:
    for name in names:
        value = os.getenv(name)
        if value and value.strip():
            return value.strip()
    return ""


def _env_bool(name: str, default: bool = False) -> bool:
    value = os.getenv(name)
    if value is None:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


def _clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "")
    return text.strip()


def _stable_hash(text: str) -> str:
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()


@dataclass
class RagChunk:
    text: str
    source: str
    source_path: str
    source_type: str
    chunk_id: str
    title: str = ""
    slide_number: int | None = None
    is_ai_llm_related: bool = False


class LocalRagStore:
    def __init__(self) -> None:
        self.debug = _env_bool("RAG_DEBUG", False)
        self.force_rebuild = _env_bool("RAG_FORCE_REBUILD", False)
        self.index_dir = Path(os.getenv("RAG_INDEX_DIR", DEFAULT_INDEX_DIR)).expanduser()
        self.embedding_deployment = _first_env(
            "AZURE_OPENAI_EMBEDDING_DEPLOYMENT",
            "AZUREAI_EMBEDDING_DEPLOYMENT",
            "AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT",
        )
        self.api_key = _first_env("AZURE_OPENAI_API_KEY", "AZUREAI_API_KEY")
        self.endpoint = _first_env("AZURE_OPENAI_ENDPOINT", "AZUREAI_ENDPOINT").rstrip("/")
        self.api_version = _first_env("AZURE_OPENAI_API_VERSION", "AZUREAI_API_VERSION") or "2024-10-21"

        self.client: AzureOpenAI | None = None
        self.index: Any | None = None
        self.chunks: list[RagChunk] = []
        self.enabled = bool(self.embedding_deployment and self.api_key and self.endpoint and faiss is not None)

        if self.debug:
            print(f"[RAG] enabled={self.enabled}")
            print(f"[RAG] embedding_deployment={self.embedding_deployment!r}")
            print(f"[RAG] endpoint_configured={bool(self.endpoint)} api_key_configured={bool(self.api_key)}")
            print(f"[RAG] faiss_available={faiss is not None}")

    def _client(self) -> AzureOpenAI:
        if self.client is None:
            self.client = AzureOpenAI(
                api_key=self.api_key,
                azure_endpoint=self.endpoint,
                api_version=self.api_version,
            )
        return self.client

    def get_source_paths(self) -> list[Path]:
        raw = os.getenv("RAG_SOURCE_DIRS") or os.getenv("RAG_SOURCE_DIR") or "knowledge"
        print(f"[RAG] Raw source setting: {raw!r}")
        paths: list[Path] = []

        for item in raw.split(";"):
            item = item.strip().strip('"').strip("'")
            if not item:
                continue
            path = Path(item).expanduser()
            # If relative, keep relative to current working directory, but resolve for logs later.
            paths.append(path)
            print(
                f"[RAG] Source candidate: {path} "
                f"exists={path.exists()} is_dir={path.is_dir()} is_file={path.is_file()}"
            )

        return paths

    def iter_source_files(self) -> list[Path]:
        files: list[Path] = []
        for source_path in self.get_source_paths():
            if not source_path.exists():
                print(f"[RAG] Missing source path: {source_path}")
                continue

            if source_path.is_file():
                if source_path.suffix.lower() in ALLOWED_EXTENSIONS:
                    files.append(source_path)
                else:
                    print(f"[RAG] Skipping unsupported file: {source_path}")
                continue

            for file_path in source_path.rglob("*"):
                if not file_path.is_file():
                    continue
                if file_path.suffix.lower() not in ALLOWED_EXTENSIONS:
                    continue
                # Skip common Office temp files
                if file_path.name.startswith("~$"):
                    continue
                files.append(file_path)

        # De-dupe while preserving order
        seen: set[str] = set()
        unique_files: list[Path] = []
        for file_path in files:
            key = str(file_path.resolve()) if file_path.exists() else str(file_path)
            if key in seen:
                continue
            seen.add(key)
            unique_files.append(file_path)

        print(f"[RAG] Found {len(unique_files)} source files")
        for file_path in unique_files[:100]:
            print(f"[RAG] FILE: {file_path}")
        if len(unique_files) > 100:
            print(f"[RAG] ...and {len(unique_files) - 100} more files")
        return unique_files

    def _is_ai_llm_related(self, text: str, file_path: Path) -> bool:
        haystack = f"{file_path.name} {text}".lower()
        terms = [
            "ai", "artificial intelligence", "llm", "large language model", "genai", "generative ai",
            "agent", "agentic", "rag", "retrieval augmented", "copilot", "chatbot", "openai",
            "machine learning", "ml", "nlp", "prompt", "semantic", "embedding",
        ]
        return any(term in haystack for term in terms)

    def _parse_file(self, file_path: Path) -> list[RagChunk]:
        suffix = file_path.suffix.lower()
        try:
            if suffix in {".html", ".htm"}:
                text = self._parse_html(file_path)
                return self._chunk_text(text, file_path, "html")
            if suffix == ".pptx":
                return self._parse_pptx(file_path)
            if suffix == ".json":
                text = self._parse_json(file_path)
                return self._chunk_text(text, file_path, "json")
            if suffix in {".txt", ".md"}:
                text = file_path.read_text(encoding="utf-8", errors="ignore")
                return self._chunk_text(text, file_path, suffix.lstrip("."))
        except Exception as exc:
            print(f"[RAG] Failed to parse {file_path}: {exc}")
        return []

    def _parse_html(self, file_path: Path) -> str:
        raw = file_path.read_text(encoding="utf-8", errors="ignore")
        if BeautifulSoup is not None:
            soup = BeautifulSoup(raw, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            return _clean_text(soup.get_text(" "))
        # crude fallback
        raw = re.sub(r"<script.*?</script>", " ", raw, flags=re.DOTALL | re.IGNORECASE)
        raw = re.sub(r"<style.*?</style>", " ", raw, flags=re.DOTALL | re.IGNORECASE)
        raw = re.sub(r"<[^>]+>", " ", raw)
        return _clean_text(raw)

    def _parse_json(self, file_path: Path) -> str:
        raw = file_path.read_text(encoding="utf-8", errors="ignore")
        try:
            data = json.loads(raw)
        except Exception:
            return raw

        # Special-ish handling for eval_Results_table.json and other structured files.
        if isinstance(data, list):
            rows = []
            for idx, item in enumerate(data, start=1):
                if isinstance(item, dict):
                    bits = [f"{k}: {v}" for k, v in item.items() if v not in (None, "")]
                    rows.append(f"Record {idx}. " + "; ".join(bits))
                else:
                    rows.append(str(item))
            return "\n".join(rows)
        if isinstance(data, dict):
            bits = []
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    bits.append(f"{key}: {json.dumps(value, ensure_ascii=False)}")
                else:
                    bits.append(f"{key}: {value}")
            return "\n".join(bits)
        return str(data)

    def _parse_pptx(self, file_path: Path) -> list[RagChunk]:
        if Presentation is None:
            print(f"[RAG] python-pptx not installed; skipping {file_path}")
            return []

        prs = Presentation(str(file_path))
        chunks: list[RagChunk] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = _clean_text(shape.text)
                    if text:
                        texts.append(text)
                # Try table cells too
                if getattr(shape, "has_table", False):
                    try:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = _clean_text(cell.text)
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                texts.append(" | ".join(row_text))
                    except Exception:
                        pass

            slide_text = _clean_text("\n".join(texts))
            if not slide_text:
                continue
            title = texts[0][:120] if texts else file_path.stem
            chunk_texts = self._split_text(slide_text)
            for part_idx, part in enumerate(chunk_texts, start=1):
                chunk_id = _stable_hash(f"{file_path.resolve()}::{slide_idx}::{part_idx}::{part[:100]}")
                chunks.append(
                    RagChunk(
                        text=part,
                        source=file_path.name,
                        source_path=str(file_path),
                        source_type="pptx",
                        chunk_id=chunk_id,
                        title=title,
                        slide_number=slide_idx,
                        is_ai_llm_related=self._is_ai_llm_related(part, file_path),
                    )
                )
        return chunks

    def _split_text(self, text: str, chunk_size: int = 1100, overlap: int = 180) -> list[str]:
        text = _clean_text(text)
        if not text:
            return []
        if len(text) <= chunk_size:
            return [text]

        chunks: list[str] = []
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            # Try to break cleanly near sentence boundary.
            if end < len(text):
                boundary = max(text.rfind(". ", start, end), text.rfind("\n", start, end))
                if boundary > start + int(chunk_size * 0.55):
                    end = boundary + 1
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(text):
                break
            start = max(0, end - overlap)
        return chunks

    def _chunk_text(self, text: str, file_path: Path, source_type: str) -> list[RagChunk]:
        parts = self._split_text(text)
        chunks: list[RagChunk] = []
        for idx, part in enumerate(parts, start=1):
            chunk_id = _stable_hash(f"{file_path.resolve()}::{idx}::{part[:100]}")
            chunks.append(
                RagChunk(
                    text=part,
                    source=file_path.name,
                    source_path=str(file_path),
                    source_type=source_type,
                    chunk_id=chunk_id,
                    title=file_path.stem,
                    is_ai_llm_related=self._is_ai_llm_related(part, file_path),
                )
            )
        return chunks

    def _embed_texts(self, texts: list[str]) -> list[list[float]]:
        if not texts:
            return []
        if not self.enabled:
            raise RuntimeError("RAG is not enabled; missing FAISS or Azure embedding config")

        embeddings: list[list[float]] = []
        batch_size = int(os.getenv("RAG_EMBED_BATCH_SIZE", "16"))
        for start in range(0, len(texts), batch_size):
            batch = texts[start:start + batch_size]
            response = self._client().embeddings.create(
                model=self.embedding_deployment,
                input=batch,
            )
            embeddings.extend([item.embedding for item in response.data])
            if self.debug:
                print(f"[RAG] Embedded {min(start + batch_size, len(texts))}/{len(texts)} chunks")
        return embeddings

    def _index_files_signature(self, files: list[Path]) -> str:
        parts = []
        for file_path in files:
            try:
                stat = file_path.stat()
                parts.append(f"{file_path.resolve()}::{stat.st_mtime_ns}::{stat.st_size}")
            except Exception:
                parts.append(str(file_path))
        return _stable_hash("\n".join(parts))

    def load_or_build(self) -> bool:
        if not self.enabled:
            print("[RAG] Disabled: missing Azure embedding config or FAISS. Concept agent will use normal model knowledge.")
            return False

        self.index_dir.mkdir(parents=True, exist_ok=True)
        index_path = self.index_dir / "index.faiss"
        chunks_path = self.index_dir / "chunks.pkl"
        sig_path = self.index_dir / "sources.sig"

        files = self.iter_source_files()
        signature = self._index_files_signature(files)

        if not files:
            print("[RAG] No source files found across configured paths")
            return False

        can_load_existing = (
            not self.force_rebuild
            and index_path.exists()
            and chunks_path.exists()
            and sig_path.exists()
            and sig_path.read_text(encoding="utf-8", errors="ignore") == signature
        )

        if can_load_existing:
            self.index = faiss.read_index(str(index_path))
            with chunks_path.open("rb") as f:
                self.chunks = pickle.load(f)
            print(f"[RAG] Loaded existing index with {len(self.chunks)} chunks from {len(files)} files")
            return True

        chunks: list[RagChunk] = []
        for file_path in files:
            file_chunks = self._parse_file(file_path)
            print(f"[RAG] Parsed {file_path}: chunks={len(file_chunks)}")
            chunks.extend(file_chunks)

        if not chunks:
            print("[RAG] No parseable content found across configured paths")
            return False

        texts = [chunk.text for chunk in chunks]
        embeddings = self._embed_texts(texts)
        if not embeddings:
            print("[RAG] No embeddings generated")
            return False

        import numpy as np

        vectors = np.array(embeddings, dtype="float32")
        faiss.normalize_L2(vectors)
        index = faiss.IndexFlatIP(vectors.shape[1])
        index.add(vectors)

        faiss.write_index(index, str(index_path))
        with chunks_path.open("wb") as f:
            pickle.dump(chunks, f)
        sig_path.write_text(signature, encoding="utf-8")

        self.index = index
        self.chunks = chunks

        by_source: dict[str, int] = {}
        for chunk in chunks:
            by_source[chunk.source] = by_source.get(chunk.source, 0) + 1

        print(f"[RAG] Indexed {len(chunks)} chunks from {len(files)} files across configured source paths")
        for source, count in sorted(by_source.items())[:50]:
            print(f"[RAG] Chunk count: {source} -> {count}")
        if len(by_source) > 50:
            print(f"[RAG] ...and {len(by_source) - 50} more sources")
        return True

    def search(self, query: str, top_k: int = 5) -> list[dict[str, Any]]:
        if self.index is None or not self.chunks:
            built = self.load_or_build()
            if not built:
                return []

        if self.index is None or not self.chunks:
            return []

        import numpy as np

        query_embedding = self._embed_texts([query])[0]
        query_vector = np.array([query_embedding], dtype="float32")
        faiss.normalize_L2(query_vector)

        limit = min(max(top_k * 3, top_k), len(self.chunks))
        scores, indices = self.index.search(query_vector, limit)

        results: list[dict[str, Any]] = []
        seen_ids: set[str] = set()
        for score, idx in zip(scores[0], indices[0]):
            if idx < 0 or idx >= len(self.chunks):
                continue
            chunk = self.chunks[int(idx)]
            if chunk.chunk_id in seen_ids:
                continue
            seen_ids.add(chunk.chunk_id)
            results.append({
                "score": float(score),
                "text": chunk.text,
                "source": chunk.source,
                "source_path": chunk.source_path,
                "source_type": chunk.source_type,
                "title": chunk.title,
                "slide_number": chunk.slide_number,
                "is_ai_llm_related": chunk.is_ai_llm_related,
            })
            if len(results) >= top_k:
                break

        if self.debug:
            print(f"[RAG] Search query={query!r} returned {len(results)} results")
            for result in results:
                print(
                    f"[RAG] Hit score={result['score']:.3f} "
                    f"source={result['source']} slide={result.get('slide_number')} "
                    f"ai_related={result.get('is_ai_llm_related')}"
                )
        return results

    def build_context_message(self, user_message: str, max_chars: int = 6000) -> str:
        results = self.search(user_message, top_k=int(os.getenv("RAG_TOP_K", "5")))
        if not results:
            return ""

        parts = [
            "Retrieved context for this answer. Use this workshop/use-case context first. "
            "If the retrieved context is insufficient, say so briefly and then use general Azure OpenAI knowledge. "
            "When relevant, mention the source workshop/deck by name. Only name-drop use-case decks when the retrieved content is actually relevant.\n"
        ]

        for idx, item in enumerate(results, start=1):
            source = item.get("source") or "Unknown source"
            slide = item.get("slide_number")
            source_label = f"{source}, slide {slide}" if slide else source
            text = _clean_text(str(item.get("text", "")))
            if not text:
                continue
            parts.append(f"\n[Source {idx}: {source_label}]\n{text}\n")

        context = "\n".join(parts).strip()
        if len(context) > max_chars:
            context = context[:max_chars].rsplit(" ", 1)[0] + "..."
        return context


rag_store = LocalRagStore()
