"""
rag.py — local FAISS RAG store for Agent 101.

Supports multiple source folders/files via RAG_SOURCE_DIRS, e.g.
  RAG_SOURCE_DIRS=/mnt/c/Users/Aparna/OneDrive - Blend 360/Use Cases;/home/aparna/agents-101/knowledge

Supported source types:
  .html, .htm, .pptx, .json, .txt, .md

Required env:
  AZURE_OPENAI_API_KEY
  AZURE_OPENAI_ENDPOINT
  AZURE_OPENAI_EMBEDDING_DEPLOYMENT

Optional env:
  RAG_SOURCE_DIRS
  RAG_SOURCE_DIR
  RAG_INDEX_DIR=.rag_index
  RAG_FORCE_REBUILD=true
  RAG_DEBUG=true
"""

from __future__ import annotations

import json
import os
import pickle
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Iterable

from dotenv import load_dotenv
from openai import AzureOpenAI

load_dotenv()

print("[RAG] Loaded multi-source rag.py v2026-05-27")

SUPPORTED_SUFFIXES = {".html", ".htm", ".pptx", ".json", ".txt", ".md"}


def _first_env(*names: str) -> str:
    for name in names:
        value = os.getenv(name)
        if value and value.strip():
            return value.strip()
    return ""


def _truthy(value: str | None) -> bool:
    return (value or "").strip().lower() in {"1", "true", "yes", "on"}


def _clean_text(text: str) -> str:
    text = re.sub(r"<script[\s\S]*?</script>", " ", text, flags=re.I)
    text = re.sub(r"<style[\s\S]*?</style>", " ", text, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _chunk_text(text: str, chunk_size: int = 1200, overlap: int = 180) -> list[str]:
    text = _clean_text(text)
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunk = text[start:end].strip()
        if len(chunk) >= 80:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = max(0, end - overlap)
    return chunks


@dataclass
class RagChunk:
    text: str
    source_name: str
    source_path: str
    source_type: str
    chunk_index: int
    metadata: dict


class LocalRagStore:
    def __init__(self) -> None:
        self.debug = _truthy(os.getenv("RAG_DEBUG"))
        self.force_rebuild = _truthy(os.getenv("RAG_FORCE_REBUILD"))
        self.index_dir = Path(os.getenv("RAG_INDEX_DIR", ".rag_index")).expanduser()
        self.index_path = self.index_dir / "faiss.index"
        self.chunks_path = self.index_dir / "chunks.pkl"

        self.embedding_deployment = _first_env(
            "AZURE_OPENAI_EMBEDDING_DEPLOYMENT",
            "AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT",
            "AZUREAI_EMBEDDING_DEPLOYMENT",
        )
        self.api_key = _first_env("AZURE_OPENAI_API_KEY", "AZUREAI_API_KEY")
        self.endpoint = _first_env("AZURE_OPENAI_ENDPOINT", "AZUREAI_ENDPOINT").rstrip("/")
        self.api_version = _first_env("AZURE_OPENAI_API_VERSION", "AZUREAI_API_VERSION") or "2024-10-21"

        self.enabled = bool(self.api_key and self.endpoint and self.embedding_deployment)
        self._client: AzureOpenAI | None = None
        self._index = None
        self._chunks: list[RagChunk] = []

        if self.debug:
            print("[RAG] Raw RAG_SOURCE_DIRS =", repr(os.getenv("RAG_SOURCE_DIRS")))
            print("[RAG] Raw RAG_SOURCE_DIR  =", repr(os.getenv("RAG_SOURCE_DIR")))
            print("[RAG] Embedding deployment =", repr(self.embedding_deployment))
            print("[RAG] Index dir =", str(self.index_dir.resolve()))

    def _get_client(self) -> AzureOpenAI:
        if self._client is None:
            self._client = AzureOpenAI(
                api_key=self.api_key,
                azure_endpoint=self.endpoint,
                api_version=self.api_version,
            )
        return self._client

    def get_source_paths(self) -> list[Path]:
        raw = os.getenv("RAG_SOURCE_DIRS") or os.getenv("RAG_SOURCE_DIR") or "knowledge"
        paths: list[Path] = []
        for item in raw.split(";"):
            item = item.strip().strip('"').strip("'")
            if not item:
                continue
            paths.append(Path(item).expanduser())

        if self.debug:
            print("[RAG] Configured source paths:")
            for path in paths:
                print(
                    f"      - {path} | exists={path.exists()} "
                    f"is_dir={path.is_dir()} is_file={path.is_file()}"
                )
        return paths

    def iter_source_files(self) -> list[Path]:
        files: list[Path] = []
        for source_path in self.get_source_paths():
            if not source_path.exists():
                print(f"[RAG] Source path does not exist: {source_path}")
                continue

            if source_path.is_file():
                if source_path.suffix.lower() in SUPPORTED_SUFFIXES:
                    files.append(source_path)
                    if self.debug:
                        print(f"[RAG] Added source file: {source_path}")
                else:
                    print(f"[RAG] Skipping unsupported file: {source_path}")
                continue

            for file_path in source_path.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_SUFFIXES:
                    files.append(file_path)
                    if self.debug:
                        print(f"[RAG] Discovered file: {file_path}")

        print(f"[RAG] Found {len(files)} source files across configured paths")
        return files

    def _read_file_text(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        try:
            if suffix in {".html", ".htm", ".txt", ".md"}:
                return file_path.read_text(encoding="utf-8", errors="ignore")

            if suffix == ".json":
                raw = file_path.read_text(encoding="utf-8", errors="ignore")
                try:
                    data = json.loads(raw)
                    return json.dumps(data, ensure_ascii=False, indent=2)
                except Exception:
                    return raw

            if suffix == ".pptx":
                try:
                    from pptx import Presentation
                except ImportError as exc:
                    print("[RAG] python-pptx is not installed. Run: pip install python-pptx")
                    raise exc
                prs = Presentation(str(file_path))
                parts: list[str] = []
                for i, slide in enumerate(prs.slides, start=1):
                    slide_text: list[str] = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    if slide_text:
                        parts.append(f"Slide {i}: " + "\n".join(slide_text))
                return "\n\n".join(parts)
        except Exception as exc:
            print(f"[RAG] Failed to parse {file_path}: {exc}")
            return ""
        return ""

    def _load_documents(self) -> list[RagChunk]:
        chunks: list[RagChunk] = []
        for file_path in self.iter_source_files():
            text = self._read_file_text(file_path)
            file_chunks = _chunk_text(text)
            print(f"[RAG] Parsed {file_path.name}: chunks={len(file_chunks)}")
            for idx, chunk in enumerate(file_chunks):
                chunks.append(
                    RagChunk(
                        text=chunk,
                        source_name=file_path.name,
                        source_path=str(file_path),
                        source_type=file_path.suffix.lower().lstrip("."),
                        chunk_index=idx,
                        metadata={"is_powerpoint": file_path.suffix.lower() == ".pptx"},
                    )
                )
        return chunks

    def _embed(self, texts: list[str]) -> list[list[float]]:
        if not texts:
            return []
        response = self._get_client().embeddings.create(
            model=self.embedding_deployment,
            input=texts,
        )
        return [item.embedding for item in response.data]

    def _build_index(self) -> None:
        import faiss
        import numpy as np

        chunks = self._load_documents()
        if not chunks:
            print("[RAG] No content found in any configured source path")
            self._chunks = []
            self._index = None
            return

        embeddings: list[list[float]] = []
        batch_size = 32
        for start in range(0, len(chunks), batch_size):
            batch = chunks[start : start + batch_size]
            embeddings.extend(self._embed([chunk.text for chunk in batch]))
            if self.debug:
                print(f"[RAG] Embedded chunks {start + 1}-{start + len(batch)} / {len(chunks)}")

        matrix = np.array(embeddings, dtype="float32")
        faiss.normalize_L2(matrix)
        index = faiss.IndexFlatIP(matrix.shape[1])
        index.add(matrix)

        self.index_dir.mkdir(parents=True, exist_ok=True)
        faiss.write_index(index, str(self.index_path))
        with self.chunks_path.open("wb") as f:
            pickle.dump([asdict(chunk) for chunk in chunks], f)

        self._index = index
        self._chunks = chunks
        unique_sources = sorted({chunk.source_name for chunk in chunks})
        print(f"[RAG] Indexed {len(chunks)} chunks from {len(unique_sources)} source files")
        for source in unique_sources[:25]:
            print(f"[RAG] Indexed source: {source}")
        if len(unique_sources) > 25:
            print(f"[RAG] ... plus {len(unique_sources) - 25} more sources")

    def _load_index(self) -> bool:
        if self.force_rebuild:
            print("[RAG] Force rebuild requested; ignoring saved index")
            return False
        if not self.index_path.exists() or not self.chunks_path.exists():
            return False
        try:
            import faiss
            self._index = faiss.read_index(str(self.index_path))
            with self.chunks_path.open("rb") as f:
                raw_chunks = pickle.load(f)
            self._chunks = [RagChunk(**item) for item in raw_chunks]
            print(f"[RAG] Loaded existing FAISS index with {len(self._chunks)} chunks")
            return True
        except Exception as exc:
            print(f"[RAG] Failed to load existing index; rebuilding. Error: {exc}")
            return False

    def ensure_ready(self) -> bool:
        if not self.enabled:
            print("[RAG] Disabled: missing Azure OpenAI embedding config")
            return False
        if self._index is not None and self._chunks:
            return True
        if self._load_index():
            return True
        self._build_index()
        return self._index is not None and bool(self._chunks)

    def search(self, query: str, top_k: int = 5) -> list[RagChunk]:
        if not self.ensure_ready():
            return []
        import faiss
        import numpy as np

        query_vec = np.array(self._embed([query]), dtype="float32")
        faiss.normalize_L2(query_vec)
        scores, indices = self._index.search(query_vec, top_k)

        results: list[RagChunk] = []
        for score, idx in zip(scores[0], indices[0]):
            if idx < 0 or idx >= len(self._chunks):
                continue
            chunk = self._chunks[idx]
            if self.debug:
                print(f"[RAG] Retrieved score={score:.3f} source={chunk.source_name} chunk={chunk.chunk_index}")
            results.append(chunk)
        return results

    def build_context(self, query: str, top_k: int = 5, max_chars: int = 6000) -> str:
        results = self.search(query, top_k=top_k)
        if not results:
            print("[RAG] No retrieved chunks for query")
            return ""

        parts: list[str] = []
        used = 0
        for chunk in results:
            header = f"Source: {chunk.source_name} | Type: {chunk.source_type} | Path: {chunk.source_path}"
            block = f"{header}\n{chunk.text}"
            if used + len(block) > max_chars:
                remaining = max_chars - used
                if remaining <= 300:
                    break
                block = block[:remaining]
            parts.append(block)
            used += len(block)
            if used >= max_chars:
                break

        context = "\n\n---\n\n".join(parts)
        print(f"[RAG] Built retrieved context from {len(results)} chunks, chars={len(context):,}")
        return context

    def augment_user_message(self, user_message: str) -> str:
        context = self.build_context(user_message)
        if not context:
            print("[RAG] Concept agent running without retrieved context")
            return user_message

        print(f"[RAG] Injecting {len(context):,} characters of retrieved context into Concept agent")
        return (
            "Use the retrieved workshop/use-case context first. "
            "If the retrieved context is insufficient, say what is missing and then use general Azure OpenAI knowledge to fill the gap. "
            "When a relevant use-case deck appears in the retrieved context, name-drop the deck/source.\n\n"
            f"Retrieved context:\n{context}\n\n"
            f"User question:\n{user_message}"
        )


rag_store = LocalRagStore()
