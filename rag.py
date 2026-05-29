"""
rag.py — cached, metadata-aware local RAG for Agent 101

What this does:
- Reads multiple source folders/files from RAG_SOURCE_DIRS, or RAG_SOURCE_DIR, or knowledge
- Supports .html/.htm, .pptx, .pdf, .json, .txt, .md
- Extracts chunks with metadata: source, path, content_type, category, topic, difficulty, AI/LLM relevance
- Builds a local FAISS index once and saves it to .rag_index/
- On future app starts, reuses the saved index unless files/config changed
- Rebuilds only when RAG_FORCE_REBUILD=true or source fingerprints/config differ
- Provides rag_store.ensure_ready() and rag_store.build_context_message() expected by main.py

Recommended env:
  RAG_ENABLED=true
  RAG_SOURCE_DIRS=knowledge;/path/to/onedrive/use-cases
  RAG_INDEX_DIR=.rag_index
  RAG_FORCE_REBUILD=false
  RAG_DEBUG=true
  AZURE_OPENAI_EMBEDDING_DEPLOYMENT=text-embedding-3-large
"""

from __future__ import annotations

import hashlib
import json
import os
import pickle
import re
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any

from dotenv import load_dotenv
from openai import AzureOpenAI

load_dotenv()

print("[RAG] Loaded cached metadata-aware rag.py")

ALLOWED_EXTENSIONS = {".html", ".htm", ".pptx", ".pdf", ".json", ".txt", ".md"}
DEFAULT_INDEX_DIR = ".rag_index"
DEFAULT_CHUNK_SIZE = 1200
DEFAULT_CHUNK_OVERLAP = 180
MANIFEST_VERSION = "metadata-cache-v2-profile-routing"

try:
    import faiss  # type: ignore
except Exception:  # pragma: no cover
    faiss = None


def _first_env(*names: str) -> str:
    for name in names:
        value = os.getenv(name)
        if value and value.strip():
            return value.strip()
    return ""


def _truthy(value: str | None, default: bool = False) -> bool:
    if value is None or value == "":
        return default
    return value.strip().lower() in {"1", "true", "yes", "y", "on"}


def _clean_text(text: str) -> str:
    text = re.sub(r"<script[\s\S]*?</script>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<style[\s\S]*?</style>", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&amp;", "&", text)
    text = re.sub(r"&lt;", "<", text)
    text = re.sub(r"&gt;", ">", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _safe_read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return path.read_text(encoding=encoding, errors="ignore")
        except Exception:
            continue
    return ""


def _sha256_for_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for block in iter(lambda: f.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


@dataclass
class RagChunk:
    text: str
    source: str
    path: str
    chunk_id: int
    content_type: str
    category: str
    topic: str
    difficulty: str
    is_ai_llm_related: bool
    metadata: dict[str, Any]


class LocalRagStore:
    def __init__(self) -> None:
        self.enabled = _truthy(os.getenv("RAG_ENABLED"), default=True)
        self.debug = _truthy(os.getenv("RAG_DEBUG"), default=False)
        self.force_rebuild = _truthy(os.getenv("RAG_FORCE_REBUILD"), default=False)
        self.index_dir = Path(os.getenv("RAG_INDEX_DIR", DEFAULT_INDEX_DIR)).expanduser()
        self.chunk_size = int(os.getenv("RAG_CHUNK_SIZE", str(DEFAULT_CHUNK_SIZE)))
        self.chunk_overlap = int(os.getenv("RAG_CHUNK_OVERLAP", str(DEFAULT_CHUNK_OVERLAP)))
        self.embedding_deployment = _first_env(
            "AZURE_OPENAI_EMBEDDING_DEPLOYMENT",
            "AZUREAI_EMBEDDING_DEPLOYMENT",
            "AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT",
        )
        self.api_key = _first_env("AZURE_OPENAI_API_KEY", "AZUREAI_API_KEY")
        self.endpoint = _first_env("AZURE_OPENAI_ENDPOINT", "AZUREAI_ENDPOINT").rstrip("/")
        self.api_version = _first_env("AZURE_OPENAI_API_VERSION", "AZUREAI_API_VERSION") or "2024-10-21"
        self.client: AzureOpenAI | None = None
        self.index = None
        self.chunks: list[RagChunk] = []
        self.ready = False
        self.last_sources: list[dict[str, Any]] = []

        print(f"[RAG] enabled={self.enabled}")
        print(f"[RAG] embedding_deployment={self.embedding_deployment!r}")
        print(f"[RAG] endpoint_configured={bool(self.endpoint)} api_key_configured={bool(self.api_key)}")
        print(f"[RAG] faiss_available={faiss is not None}")

    # ---------- public API expected by main.py ----------

    def ensure_ready(self) -> None:
        if self.ready:
            return
        if not self.enabled:
            print("[RAG] disabled")
            return
        if faiss is None:
            print("[RAG] unavailable: faiss is not installed")
            return
        if not self.embedding_deployment or not self.api_key or not self.endpoint:
            print("[RAG] unavailable: missing Azure OpenAI embedding config")
            return

        self.index_dir.mkdir(parents=True, exist_ok=True)

        if not self.force_rebuild and self._saved_index_is_current():
            self._load_index()
            self.ready = self.index is not None and bool(self.chunks)
            if self.ready:
                print(f"[RAG] Loaded cached index with {len(self.chunks)} chunks from {self.index_dir.resolve()}")
                return

        self.build_index()

    def build_context_message(self, user_message: str, max_chars: int | None = None) -> str:
        """Build a source-attributed RAG context block for the Concept agent.

        The returned string is injected as context by main.py. It includes numbered
        source labels like [1], [2], etc. so the Concept agent can cite sources in
        the final response and append a Sources Used section.
        """
        self.ensure_ready()
        if not self.ready:
            return ""

        profile = self.infer_retrieval_profile(user_message)
        if max_chars is None:
            max_chars = profile["max_context_chars"]

        results = self.search(user_message, top_k=profile["top_k"], profile=profile)
        if not results:
            self.last_sources = []
            print("[RAG] No retrieved sources matched this query")
            return ""

        # Build stable, deduplicated source numbers.
        source_numbers: dict[str, int] = {}
        source_details: list[dict[str, Any]] = []

        for item in results:
            meta = item.get("metadata", {})
            source = meta.get("source", item.get("source", "Unknown source"))
            path = meta.get("path", "")
            key = f"{source}|{path}"

            if key not in source_numbers:
                source_numbers[key] = len(source_numbers) + 1
                source_details.append(
                    {
                        "number": source_numbers[key],
                        "source": source,
                        "path": path,
                        "category": meta.get("category", "general"),
                        "content_type": meta.get("content_type", "document"),
                        "topic": meta.get("topic", "general"),
                        "difficulty": meta.get("difficulty", "general"),
                        "is_ai_llm_related": meta.get("is_ai_llm_related", False),
                    }
                )

        self.last_sources = source_details

        print(
            f"[RAG] Retrieval profile: difficulty={profile['difficulty']} "
            f"intent={profile['intent']} topic={profile['topic']} "
            f"top_k={profile['top_k']} max_chars={max_chars}"
        )
        print("[RAG] Retrieved sources:")
        for detail in source_details:
            print(
                f"      [{detail['number']}] {detail['source']} "
                f"content_type={detail['content_type']} "
                f"category={detail['category']} topic={detail['topic']} "
                f"difficulty={detail['difficulty']} ai_llm_related={detail['is_ai_llm_related']}"
            )

        if self.debug:
            print("[RAG] Retrieved chunks:")
            for idx, item in enumerate(results, start=1):
                meta = item.get("metadata", {})
                print(
                    f"      chunk={idx} source={meta.get('source', item.get('source'))} "
                    f"score={item.get('score', 0.0):.4f} "
                    f"vector={item.get('vector_score', 0.0):.4f} "
                    f"boost={meta.get('metadata_boost', 0.0):.4f} "
                    f"chunk_id={meta.get('chunk_id')}"
                )

        source_index_text = "\n".join(
            (
                f"[{detail['number']}] {detail['source']} "
                f"(content_type={detail['content_type']}; category={detail['category']}; "
                f"topic={detail['topic']}; difficulty={detail['difficulty']})"
            )
            for detail in source_details
        )

        parts = [
            "Retrieved context for Agent 101. Use this context first before general model knowledge.\n"
            "Citation rules for the final answer:\n"
            "- When you use a retrieved fact, cite it inline with the source number, e.g. [1].\n"
            "- If a use-case deck is relevant, mention the deck/use-case name naturally.\n"
            "- If the retrieved context is insufficient, say so briefly and then use general Azure OpenAI knowledge.\n"
            "- End the response with a short 'Sources Used' section listing only sources actually used.\n"
            "- Do not invent sources, page numbers, URLs, or deck names.\n"
            f"Retrieval profile inferred from the question: difficulty={profile['difficulty']}; "
            f"intent={profile['intent']}; topic={profile['topic']}.\n\n"
            "Source index:\n"
            f"{source_index_text}\n"
        ]

        used = sum(len(part) for part in parts)

        for item in results:
            meta = item.get("metadata", {})
            source = meta.get("source", item.get("source", "Unknown source"))
            path = meta.get("path", "")
            key = f"{source}|{path}"
            source_no = source_numbers.get(key, 0)

            category = meta.get("category", "general")
            topic = meta.get("topic", "general")
            difficulty = meta.get("difficulty", "general")
            content_type = meta.get("content_type", "document")
            ai_flag = meta.get("is_ai_llm_related", False)
            chunk_id = meta.get("chunk_id", "")
            score = item.get("score", 0.0)
            vector_score = item.get("vector_score", 0.0)
            boost = meta.get("metadata_boost", 0.0)
            text = item.get("text", "").strip()
            if not text:
                continue

            block = (
                f"\n[Source {source_no}: {source}]\n"
                f"metadata: content_type={content_type}; category={category}; topic={topic}; "
                f"difficulty={difficulty}; ai_llm_related={ai_flag}; chunk_id={chunk_id}; "
                f"retrieval_score={score:.4f}; vector_score={vector_score:.4f}; metadata_boost={boost:.4f}\n"
                f"{text}\n"
            )
            if used + len(block) > max_chars:
                remaining = max_chars - used
                if remaining > 400:
                    parts.append(block[:remaining])
                    used += remaining
                break
            parts.append(block)
            used += len(block)

        context = "\n".join(parts).strip()
        if context:
            print(f"[RAG] Injecting {len(context):,} characters of source-attributed context into Concept agent")
        return context


    def get_last_sources(self) -> list[dict[str, Any]]:
        """Return the deduplicated sources from the most recent retrieval."""
        return list(self.last_sources)

    # ---------- source discovery + caching ----------

    def get_source_paths(self) -> list[Path]:
        raw = os.getenv("RAG_SOURCE_DIRS") or os.getenv("RAG_SOURCE_DIR") or "knowledge"
        if self.debug:
            print(f"[RAG] Raw source setting: {raw!r}")

        paths: list[Path] = []
        for item in raw.split(";"):
            item = item.strip()
            if not item:
                continue
            path = Path(item).expanduser()
            paths.append(path)
            if self.debug:
                print(
                    f"[RAG] Source candidate: {path} "
                    f"exists={path.exists()} is_dir={path.is_dir()} is_file={path.is_file()}"
                )
        return paths

    def iter_source_files(self) -> list[Path]:
        files: list[Path] = []
        seen: set[str] = set()

        for source_path in self.get_source_paths():
            if not source_path.exists():
                print(f"[RAG] Missing source path: {source_path}")
                continue

            candidates = [source_path] if source_path.is_file() else list(source_path.rglob("*"))
            for file_path in candidates:
                if not file_path.is_file():
                    continue
                if file_path.suffix.lower() not in ALLOWED_EXTENSIONS:
                    continue
                key = str(file_path.resolve())
                if key in seen:
                    continue
                seen.add(key)
                files.append(file_path)

        files.sort(key=lambda p: str(p).lower())
        if self.debug:
            print(f"[RAG] Found {len(files)} supported source files")
            for file_path in files[:100]:
                print(f"[RAG] FILE: {file_path}")
        return files

    def _current_manifest(self) -> dict[str, Any]:
        files = self.iter_source_files()
        file_records = []
        for path in files:
            try:
                stat = path.stat()
                file_records.append(
                    {
                        "path": str(path.resolve()),
                        "size": stat.st_size,
                        "mtime_ns": stat.st_mtime_ns,
                        # Hash is slower but robust; can disable once stable if desired.
                        "sha256": _sha256_for_file(path),
                    }
                )
            except Exception as exc:
                print(f"[RAG] Could not fingerprint {path}: {exc}")

        return {
            "version": MANIFEST_VERSION,
            "embedding_deployment": self.embedding_deployment,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "sources_raw": os.getenv("RAG_SOURCE_DIRS") or os.getenv("RAG_SOURCE_DIR") or "knowledge",
            "files": file_records,
        }

    def _manifest_path(self) -> Path:
        return self.index_dir / "manifest.json"

    def _chunks_path(self) -> Path:
        return self.index_dir / "chunks.pkl"

    def _faiss_path(self) -> Path:
        return self.index_dir / "index.faiss"

    def _saved_index_is_current(self) -> bool:
        if not self._manifest_path().exists() or not self._chunks_path().exists() or not self._faiss_path().exists():
            if self.debug:
                print("[RAG] Cached index missing one or more files; rebuild needed")
            return False

        try:
            saved = json.loads(self._manifest_path().read_text(encoding="utf-8"))
            current = self._current_manifest()
        except Exception as exc:
            print(f"[RAG] Could not compare manifest; rebuild needed: {exc}")
            return False

        is_current = saved == current
        if self.debug:
            print(f"[RAG] Cached index current={is_current}")
            if not is_current:
                print("[RAG] Source/config changed; rebuild needed")
        return is_current

    # ---------- parsing/chunking ----------

    def parse_file(self, path: Path) -> str:
        suffix = path.suffix.lower()
        try:
            if suffix in {".html", ".htm"}:
                return _clean_text(_safe_read_text(path))
            if suffix in {".txt", ".md"}:
                return _clean_text(_safe_read_text(path))
            if suffix == ".json":
                return self._parse_json(path)
            if suffix == ".pptx":
                return self._parse_pptx(path)
            if suffix == ".pdf":
                return self._parse_pdf(path)
        except Exception as exc:
            print(f"[RAG] Failed parsing {path}: {exc}")
        return ""

    def _parse_json(self, path: Path) -> str:
        raw = _safe_read_text(path)
        if not raw.strip():
            return ""
        try:
            data = json.loads(raw)
        except Exception:
            return _clean_text(raw)
        return _clean_text(json.dumps(data, ensure_ascii=False, indent=2))

    def _parse_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            print("[RAG] python-pptx is not installed; skipping PPTX parsing")
            return ""

        prs = Presentation(str(path))
        slide_texts: list[str] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    cleaned = _clean_text(shape.text)
                    if cleaned:
                        texts.append(cleaned)
            if texts:
                slide_texts.append(f"[Slide {slide_idx}]\n" + "\n".join(texts))
        return "\n\n".join(slide_texts)

    def _parse_pdf(self, path: Path) -> str:
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception:
            print("[RAG] pypdf is not installed; skipping PDF parsing")
            return ""

        reader = PdfReader(str(path))
        pages: list[str] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            text = _clean_text(text)
            if text:
                pages.append(f"[Page {i}]\n{text}")
        return "\n\n".join(pages)

    def chunk_text(self, text: str) -> list[str]:
        text = _clean_text(text)
        if not text:
            return []

        chunks: list[str] = []
        start = 0
        n = len(text)
        while start < n:
            end = min(start + self.chunk_size, n)
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= n:
                break
            start = max(0, end - self.chunk_overlap)
        return chunks

    def build_chunks(self) -> list[RagChunk]:
        chunks: list[RagChunk] = []
        for file_path in self.iter_source_files():
            text = self.parse_file(file_path)
            file_chunks = self.chunk_text(text)
            if self.debug:
                print(f"[RAG] Parsed {file_path.name}: chars={len(text):,} chunks={len(file_chunks)}")
            for idx, chunk_text in enumerate(file_chunks):
                metadata = self.infer_metadata(file_path, chunk_text)
                chunks.append(
                    RagChunk(
                        text=chunk_text,
                        source=file_path.name,
                        path=str(file_path),
                        chunk_id=idx,
                        content_type=metadata["content_type"],
                        category=metadata["category"],
                        topic=metadata["topic"],
                        difficulty=metadata["difficulty"],
                        is_ai_llm_related=metadata["is_ai_llm_related"],
                        metadata=metadata,
                    )
                )
        return chunks

    # ---------- metadata-aware retrieval ----------

    def infer_metadata(self, path: Path, text: str) -> dict[str, Any]:
        lower_path = str(path).lower()
        lower_text = text.lower()
        suffix = path.suffix.lower()

        if suffix == ".pptx":
            content_type = "deck"
        elif suffix == ".pdf":
            content_type = "paper_or_pdf"
        elif suffix in {".html", ".htm"}:
            content_type = "workshop_html"
        elif suffix == ".json":
            content_type = "structured_json"
        else:
            content_type = "text"

        path_parts = [part.lower() for part in path.parts]
        category = "general"
        for candidate in (
            "foundations",
            "guides",
            "mcp",
            "production",
            "use-cases",
            "use_cases",
            "workshop",
            "transcripts",
            "advanced-rag",
            "advanced_rag",
            "knowledge",
        ):
            if candidate in path_parts or candidate in lower_path:
                category = candidate.replace("_", "-")
                break
        if content_type == "deck" and category == "general":
            category = "use-cases"
        if content_type == "workshop_html":
            category = "workshop"

        topic_keywords = {
            "agents": ["agent", "agentic", "autonomous", "planning", "act", "observe"],
            "rag": ["rag", "retrieval", "vector", "embedding", "chunk", "faiss", "search"],
            "tool-calling": ["tool calling", "function calling", "tools", "api", "mcp", "model context protocol"],
            "multi-agent": ["multi-agent", "multi agent", "orchestrator", "handoff", "coordination"],
            "memory": ["memory", "reflection", "state", "context", "session"],
            "evals": ["eval", "evaluation", "benchmark", "judge", "regression", "golden dataset"],
            "production": ["production", "observability", "tracing", "latency", "cost", "guardrail", "deploy"],
            "healthcare": ["medical", "health", "clinical", "patient", "doctor", "diagnosis", "hospital"],
            "retail": ["retail", "customer", "promotion", "pricing", "merchandising", "commerce"],
            "finance": ["finance", "financial", "bank", "risk", "fraud", "insurance"],
        }
        topic_scores = {
            topic: sum(1 for kw in keywords if kw in lower_text or kw in lower_path)
            for topic, keywords in topic_keywords.items()
        }
        topic = max(topic_scores, key=topic_scores.get)
        if topic_scores[topic] == 0:
            topic = "general"

        if any(term in lower_text for term in ["advanced", "graph", "self-rag", "raptor", "reflection", "multi-agent"]):
            difficulty = "advanced"
        elif any(term in lower_text for term in ["architecture", "production", "deploy", "evaluation", "observability"]):
            difficulty = "intermediate"
        elif any(term in lower_text for term in ["101", "intro", "foundations", "beginner", "what is"]):
            difficulty = "beginner"
        else:
            difficulty = "general"

        ai_terms = [
            "ai", "llm", "large language model", "agent", "agentic", "rag", "embedding",
            "prompt", "mcp", "tool calling", "function calling", "generative", "openai",
            "claude", "langchain", "langgraph", "autogen", "vector", "retrieval",
        ]
        is_ai_llm_related = any(term in lower_text or term in lower_path for term in ai_terms)

        return {
            "source": path.name,
            "path": str(path),
            "content_type": content_type,
            "category": category,
            "topic": topic,
            "difficulty": difficulty,
            "is_ai_llm_related": is_ai_llm_related,
        }

    def infer_retrieval_profile(self, query: str) -> dict[str, Any]:
        """Infer retrieval difficulty/intent/topic from the question itself.

        No user-level override is needed: the wording of the question controls how much
        context we retrieve and which metadata we prefer. This saves tokens for simple
        questions and allows deeper context for advanced/production questions.
        """
        lower = query.lower()

        topic = "general"
        topic_rules = [
            ("healthcare", ["medical", "health", "clinical", "patient", "doctor", "hospital", "diagnosis"]),
            ("retail", ["retail", "customer", "promotion", "pricing", "merchandising", "commerce"]),
            ("finance", ["finance", "financial", "bank", "risk", "fraud", "insurance"]),
            ("rag", ["rag", "retrieval", "embedding", "vector", "chunk", "faiss", "knowledge base"]),
            ("tool-calling", ["tool", "function", "mcp", "api", "integration", "connector"]),
            ("multi-agent", ["multi-agent", "multi agent", "orchestr", "handoff", "swarm"]),
            ("memory", ["memory", "context engineering", "context", "state", "session"]),
            ("evals", ["eval", "evaluation", "test", "judge", "benchmark", "regression"]),
            ("production", ["deploy", "production", "monitor", "observability", "latency", "cost", "guardrail"]),
            ("agents", ["agent", "agentic", "llm", "large language model"]),
        ]
        for candidate, keywords in topic_rules:
            if any(keyword in lower for keyword in keywords):
                topic = candidate
                break

        if any(x in lower for x in ["code", "build", "implement", "python", "script", "function", "example code"]):
            intent = "implementation"
        elif any(x in lower for x in ["deploy", "production", "monitor", "eval", "cost", "latency", "security", "guardrail"]):
            intent = "production"
        elif any(x in lower for x in ["use case", "case study", "client", "industry", "example", "deck"]):
            intent = "use_case"
        else:
            intent = "concept"

        beginner_signals = [
            "what is", "what's", "explain", "simple", "simply", "beginner", "101",
            "basics", "difference between", "how does", "how do", "intro"
        ]
        advanced_signals = [
            "advanced", "architecture", "tradeoff", "trade-off", "deep dive", "scale",
            "production", "deploy", "eval", "observability", "latency", "security",
            "failure mode", "benchmark", "optimize", "cost", "multi-agent", "graph", "self-rag"
        ]
        implementation_signals = ["code", "implement", "build", "python", "api", "class", "function"]

        if any(x in lower for x in advanced_signals):
            difficulty = "advanced"
        elif any(x in lower for x in implementation_signals):
            difficulty = "intermediate"
        elif any(x in lower for x in beginner_signals):
            difficulty = "beginner"
        else:
            difficulty = "intermediate"

        if difficulty == "beginner":
            top_k = 3
            max_context_chars = 2500
        elif difficulty == "advanced":
            top_k = 7
            max_context_chars = 6500
        else:
            top_k = 5
            max_context_chars = 4500

        return {
            "topic": topic,
            "intent": intent,
            "difficulty": difficulty,
            "top_k": top_k,
            "max_context_chars": max_context_chars,
        }

    def metadata_boost(self, chunk: RagChunk, query: str, profile: dict[str, Any] | None = None) -> float:
        profile = profile or self.infer_retrieval_profile(query)
        boost = 0.0
        text_lower = chunk.text.lower()

        if chunk.topic == profile["topic"]:
            boost += 0.14
        elif profile["topic"] != "general" and profile["topic"] in text_lower:
            boost += 0.06

        if chunk.difficulty == profile["difficulty"]:
            boost += 0.10
        elif chunk.difficulty == "general":
            boost += 0.02
        elif profile["difficulty"] == "beginner" and chunk.difficulty == "advanced":
            boost -= 0.10

        intent = profile["intent"]
        if intent == "concept":
            if chunk.content_type in {"workshop_html", "text", "paper_or_pdf"}:
                boost += 0.08
            if chunk.category in {"workshop", "foundations", "guides", "knowledge"}:
                boost += 0.08
            if chunk.category == "use-cases":
                boost -= 0.04
        elif intent == "implementation":
            if chunk.category in {"guides", "production", "mcp", "advanced-rag", "knowledge"}:
                boost += 0.08
            if any(term in text_lower for term in ["python", "code", "api", "implementation", "function"]):
                boost += 0.08
        elif intent == "production":
            if chunk.category == "production" or chunk.topic in {"production", "evals"}:
                boost += 0.16
            if any(term in text_lower for term in ["monitor", "trace", "latency", "cost", "eval", "guardrail"]):
                boost += 0.08
        elif intent == "use_case":
            if chunk.category == "use-cases" or chunk.content_type == "deck":
                boost += 0.18
            if chunk.is_ai_llm_related:
                boost += 0.08

        if chunk.is_ai_llm_related:
            boost += 0.06
        elif chunk.category == "use-cases":
            boost -= 0.10

        return boost

    # ---------- FAISS embedding/index ----------

    def _get_client(self) -> AzureOpenAI:
        if self.client is None:
            self.client = AzureOpenAI(
                api_key=self.api_key,
                azure_endpoint=self.endpoint,
                api_version=self.api_version,
            )
        return self.client

    def embed_texts(self, texts: list[str]) -> list[list[float]]:
        if not texts:
            return []
        client = self._get_client()
        batch_size = int(os.getenv("RAG_EMBED_BATCH_SIZE", "16"))
        embeddings: list[list[float]] = []
        for start in range(0, len(texts), batch_size):
            batch = texts[start:start + batch_size]
            response = client.embeddings.create(model=self.embedding_deployment, input=batch)
            embeddings.extend([item.embedding for item in response.data])
            if self.debug:
                print(f"[RAG] Embedded {min(start + batch_size, len(texts))}/{len(texts)} chunks")
        return embeddings

    def build_index(self) -> None:
        if faiss is None:
            print("[RAG] Cannot build index: faiss unavailable")
            return

        manifest = self._current_manifest()
        chunks = self.build_chunks()
        if not chunks:
            print("[RAG] No chunks built from configured sources")
            return

        embeddings = self.embed_texts([chunk.text for chunk in chunks])
        if not embeddings:
            print("[RAG] No embeddings created")
            return

        import numpy as np

        vectors = np.array(embeddings, dtype="float32")
        faiss.normalize_L2(vectors)
        index = faiss.IndexFlatIP(vectors.shape[1])
        index.add(vectors)

        self.index = index
        self.chunks = chunks
        self.ready = True

        self.index_dir.mkdir(parents=True, exist_ok=True)
        faiss.write_index(index, str(self._faiss_path()))
        with self._chunks_path().open("wb") as f:
            pickle.dump([asdict(chunk) for chunk in chunks], f)
        self._manifest_path().write_text(json.dumps(manifest, indent=2), encoding="utf-8")

        source_count = len({chunk.path for chunk in chunks})
        print(f"[RAG] Built and cached index: {len(chunks)} chunks from {source_count} source files")

    def _load_index(self) -> None:
        if faiss is None:
            return
        self.index = faiss.read_index(str(self._faiss_path()))
        raw_chunks = pickle.loads(self._chunks_path().read_bytes())
        self.chunks = [RagChunk(**item) for item in raw_chunks]

    def search(self, query: str, top_k: int = 7, profile: dict[str, Any] | None = None) -> list[dict[str, Any]]:
        if not self.ready or self.index is None or not self.chunks:
            return []

        import numpy as np

        query_embedding = self.embed_texts([query])[0]
        q = np.array([query_embedding], dtype="float32")
        faiss.normalize_L2(q)

        profile = profile or self.infer_retrieval_profile(query)
        # Over-fetch, then metadata-rerank.
        fetch_k = min(max(top_k * 5, top_k), len(self.chunks))
        scores, indices = self.index.search(q, fetch_k)

        ranked: list[dict[str, Any]] = []
        for score, idx in zip(scores[0].tolist(), indices[0].tolist()):
            if idx < 0 or idx >= len(self.chunks):
                continue
            chunk = self.chunks[idx]
            boost = self.metadata_boost(chunk, query, profile)
            final_score = float(score) + boost
            ranked.append(
                {
                    "text": chunk.text,
                    "source": chunk.source,
                    "score": final_score,
                    "vector_score": float(score),
                    "metadata": {
                        **chunk.metadata,
                        "chunk_id": chunk.chunk_id,
                        "vector_score": float(score),
                        "metadata_boost": boost,
                    },
                }
            )

        ranked.sort(key=lambda item: item["score"], reverse=True)

        # Diversity: avoid filling all slots from the same file unless highly relevant.
        selected: list[dict[str, Any]] = []
        per_source_count: dict[str, int] = {}
        max_per_source = int(os.getenv("RAG_MAX_CHUNKS_PER_SOURCE", "3"))
        min_score = float(os.getenv("RAG_MIN_SCORE", "0.20"))

        for item in ranked:
            if item["score"] < min_score:
                continue
            source = item["metadata"].get("source", item["source"])
            if per_source_count.get(source, 0) >= max_per_source:
                continue
            selected.append(item)
            per_source_count[source] = per_source_count.get(source, 0) + 1
            if len(selected) >= top_k:
                break

        if self.debug:
            print(f"[RAG] Search query={query!r} profile={profile}")
            for item in selected:
                meta = item["metadata"]
                print(
                    f"[RAG] HIT score={item['score']:.4f} vector={item['vector_score']:.4f} "
                    f"boost={meta.get('metadata_boost', 0):.4f} source={meta.get('source')} "
                    f"category={meta.get('category')} topic={meta.get('topic')} ai={meta.get('is_ai_llm_related')}"
                )

        return selected


rag_store = LocalRagStore()
